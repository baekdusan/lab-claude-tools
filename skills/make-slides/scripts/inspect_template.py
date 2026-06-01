#!/usr/bin/env python3
"""템플릿 .pptx의 레이아웃과 플레이스홀더 구조를 출력한다.

빌드 스크립트를 작성하기 전에 반드시 실행해서, 각 레이아웃의 인덱스와
그 안의 플레이스홀더 idx / type / name을 파악한다. python-pptx로 내용을
채우려면 이 인덱스를 알아야 한다.

사용법:
    python3 inspect_template.py [템플릿.pptx]
    # 인자 생략 시 번들된 기본 템플릿(../assets/template.pptx)을 스캔
"""
import os
import sys

try:
    from pptx import Presentation
except ImportError:
    sys.exit("python-pptx가 없습니다. `pip3 install python-pptx` 후 다시 실행하세요.")

# 스킬 디렉토리에 번들된 기본 템플릿 경로 (이 스크립트 기준 ../assets/template.pptx)
DEFAULT_TEMPLATE = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "..", "assets", "template.pptx"
)


def main(path: str) -> None:
    prs = Presentation(path)
    print(f"# 템플릿: {path}")
    print(f"슬라이드 크기: {prs.slide_width} x {prs.slide_height} "
          f"(EMU; /914400 = inch)\n")

    print(f"## 슬라이드 레이아웃 ({len(prs.slide_layouts)}개)")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"\n[layout {i}] \"{layout.name}\"")
        for ph in layout.placeholders:
            print(f"    - idx={ph.placeholder_format.idx:<3} "
                  f"type={str(ph.placeholder_format.type):<18} "
                  f"name=\"{ph.name}\"")

    # 템플릿에 이미 들어있는 예시 슬라이드(있으면) 구조도 표시
    if len(prs.slides):
        print(f"\n## 템플릿에 포함된 기존 슬라이드 ({len(prs.slides)}개)")
        for i, slide in enumerate(prs.slides):
            ln = slide.slide_layout.name
            print(f"\n[slide {i}] layout=\"{ln}\"")
            for ph in slide.placeholders:
                txt = (ph.text[:40] + "…") if getattr(ph, "text", "") and len(ph.text) > 40 else getattr(ph, "text", "")
                print(f"    - idx={ph.placeholder_format.idx:<3} "
                      f"name=\"{ph.name}\"  text=\"{txt}\"")


if __name__ == "__main__":
    path = sys.argv[1] if len(sys.argv) > 1 else DEFAULT_TEMPLATE
    if not os.path.exists(path):
        sys.exit(f"템플릿을 찾을 수 없습니다: {path}")
    main(path)
