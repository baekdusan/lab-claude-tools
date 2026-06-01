#!/usr/bin/env python3
"""빌드 스크립트 출발점(예시). 실제 작업 시 이 파일을 복사해 수정한다.

전제: inspect_template.py로 레이아웃/플레이스홀더 idx를 미리 확인했다.
아래 LAYOUT_*, PH_* 상수를 그 결과에 맞게 반드시 교체할 것.
"""
import os
from pptx import Presentation

# 기본값: 스킬에 번들된 공용 템플릿. 사용자가 다른 템플릿을 주면 이 값만 바꾼다.
TEMPLATE = os.path.join(
    os.path.dirname(os.path.abspath(__file__)), "..", "assets", "template.pptx"
)
OUTPUT   = "발표자료.pptx"    # ← 출력 경로

# --- inspect 결과를 보고 교체할 값들 -------------------------------------
LAYOUT_TITLE = 0   # 타이틀 슬라이드 레이아웃 인덱스
LAYOUT_BODY  = 1   # 제목+본문 레이아웃 인덱스
PH_TITLE     = 0   # 제목 플레이스홀더 idx
PH_SUBTITLE  = 1   # 부제 플레이스홀더 idx (타이틀 레이아웃)
PH_BODY      = 1   # 본문 플레이스홀더 idx (본문 레이아웃)
# ------------------------------------------------------------------------


def clear_slides(prs):
    """템플릿에 포함된 기존 슬라이드를 모두 제거(마스터·레이아웃은 보존).

    python-pptx엔 슬라이드 삭제 공식 API가 없어 sldIdLst를 직접 비운다.
    목록 참조뿐 아니라 관계(rel)도 끊어야 좀비 파트가 남지 않는다
    (안 그러면 저장 시 'Duplicate name' 경고와 함께 패키지가 지저분해짐).
    레이아웃만 빌려 본문을 처음부터 새로 채우는 방식이므로 빌드 맨 앞에서 호출.
    """
    RID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    xml_slides = prs.slides._sldIdLst
    for sld_id in list(xml_slides):
        prs.part.drop_rel(sld_id.get(RID))
        xml_slides.remove(sld_id)


def add_title(prs, title, subtitle):
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    s.placeholders[PH_TITLE].text = title
    s.placeholders[PH_SUBTITLE].text = subtitle
    return s


def add_bullets(prs, title, bullets, notes=None):
    """bullets: [(텍스트, level), ...] 또는 [텍스트, ...]"""
    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BODY])
    s.placeholders[PH_TITLE].text = title
    tf = s.placeholders[PH_BODY].text_frame
    for i, item in enumerate(bullets):
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = level
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def main():
    prs = Presentation(TEMPLATE)
    clear_slides(prs)   # 레이아웃만 빌리고 본문은 새로

    add_title(prs, "발표 제목", "한 줄 부제 · 발표자 · 2026")

    add_bullets(
        prs,
        "배경 — 왜 이 문제인가",
        ["현재 분야는 ~한 상황 (Context)",
         "그런데 ~가 해결되지 않음 (Problem)",
         ("안 풀면 ~한 손해 (Cost: So what?)", 1)],
        notes="청중에게 익숙한 사례로 시작",
    )

    add_bullets(
        prs,
        "제안 기법은 baseline 대비 +2.3점",
        ["방법 포인트 1", "방법 포인트 2", "방법 포인트 3"],
    )

    prs.save(OUTPUT)
    print(f"저장됨: {OUTPUT}  (슬라이드 {len(prs.slides)}장)")


if __name__ == "__main__":
    main()
